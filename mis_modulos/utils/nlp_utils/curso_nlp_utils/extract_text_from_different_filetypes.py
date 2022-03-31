"""
Create a python function that is responsible for extracting the text of an HTML document to be processed by NLP
"""
import xml

import PyPDF2 as PyPDF2
import docx as docx
import pptx
import pytesseract as pytesseract
from bs4 import BeautifulSoup
from wasabi import markdown


def extract_html_text(html_text):
    extracted_text = ''
    soup = BeautifulSoup(html_text, 'html.parser')
    for text in soup.find_all('p'):
        extracted_text += text.get_text()
    return extracted_text


def extract_pdf_text(pdf_file):
    extracted_text = ''
    pdf_file = open(pdf_file, 'rb')
    pdf_reader = PyPDF2.PdfFileReader(pdf_file)
    for page in range(pdf_reader.numPages):
        extracted_text += pdf_reader.getPage(page).extractText()
    return extracted_text


def extract_image_pdf_text(pdf_file):
    extracted_text = ''
    pdf_file = open(pdf_file, 'rb')
    pdf_reader = PyPDF2.PdfFileReader(pdf_file)
    for page in range(pdf_reader.numPages):
        page_obj = pdf_reader.getPage(page)
        extracted_text += pytesseract.image_to_string(page_obj.extractText())
    return extracted_text


def extract_image_doc_text(doc_file):
    extracted_text = ''
    doc_file = open(doc_file, 'rb')
    doc_reader = docx.Document(doc_file)
    for paragraph in doc_reader.paragraphs:
        extracted_text += paragraph.text
    return extracted_text


def extract_xml_text(xml_file):
    extracted_text = ''
    xml_file = open(xml_file, 'rb')
    xml_reader = xml.dom.minidom.parse(xml_file)
    for node in xml_reader.getElementsByTagName('p'):
        extracted_text += node.firstChild.nodeValue
    return extracted_text


def extract_ppt_text(ppt_file):
    extracted_text = ''
    ppt_file = open(ppt_file, 'rb')
    ppt_reader = pptx.Presentation(ppt_file)
    for slide in ppt_reader.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    extracted_text += run.text
    return extracted_text


def extract_md_text(md_file):
    extracted_text = ''
    md_file = open(md_file, 'rb')
    md_reader = markdown.Markdown()
    extracted_text = md_reader.convert(md_file.read())
    return extracted_text
